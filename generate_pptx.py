from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AryaX ASI"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(60)
    
    subtitle.text = "Level 3000 Galactic Artificial Super Intelligence\nPowered by Advanced Multi-Model Consensus"
    
    # 2. Vision Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "The Vision: Beyond a Chatbot"
    tf = body_shape.text_frame
    tf.text = "AryaX is designed as a complete Autonomous Ecosystem."
    
    p = tf.add_paragraph()
    p.text = "🧠 Multi-Model Consensus: Merges Gemini Pro, Flash, and GPT-4o for supreme reasoning."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "🌍 Real-Time Awareness: Connects to live Crypto, Stocks, News, and Weather."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "🤖 Autonomous Agent: Can research, analyze, and execute tasks without human help."
    p.level = 1

    # 3. Technical Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Core Technical Architecture"
    tf = body_shape.text_frame
    tf.text = "Built on a robust, scalable backend."
    
    p = tf.add_paragraph()
    p.text = "Backend: Python Flask & Gunicorn for high performance."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Memory: Pinecone Vector DB for infinite, long-term context recall."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Execution: Judge0 Sandbox for secure, multi-language code execution."
    p.level = 1

    # 4. Multimodal & Voice
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "True Multimodal Experience"
    tf = body_shape.text_frame
    tf.text = "Interacting with AryaX feels human."
    
    p = tf.add_paragraph()
    p.text = "🎙️ Neural Listening: OpenAI Whisper API for flawless voice recognition."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "🔊 Realistic Speech: ElevenLabs API for human-like emotional voice output."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "🎨 Generation: Image & File (PDF/CSV/Excel) generation capabilities on the fly."
    p.level = 1

    # 5. The Future
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "The Future of AryaX"
    tf = body_shape.text_frame
    tf.text = "AryaX is ready for deployment and scaling."
    
    p = tf.add_paragraph()
    p.text = "• Deploying to Render with auto-scaling instances."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Ready for enterprise integration and B2B solutions."
    p.level = 1

    prs.save("AryaX_Pitch_Deck.pptx")
    print("Presentation saved as AryaX_Pitch_Deck.pptx")

if __name__ == '__main__':
    create_presentation()

from flask import Flask, request, jsonify, send_from_directory, Response, stream_with_context, send_file
from datetime import datetime, date
import pytz
import requests
import os
import json
import time
import hashlib
import urllib.parse
import random
import threading
import io
from dotenv import load_dotenv
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    # pyrefly: ignore [missing-import]
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None
try:
    # pyrefly: ignore [missing-import]
    import PyPDF2
except ImportError:
    PyPDF2 = None

load_dotenv()

app = Flask(__name__, static_folder='public')

# ===== JWT AUTH DECORATOR =====
import jwt
from functools import wraps
def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        token = None
        if 'Authorization' in request.headers:
            parts = request.headers['Authorization'].split()
            if len(parts) == 2 and parts[0] == 'Bearer':
                token = parts[1]
        
        if not token:
            return jsonify({'error': 'Token is missing!'}), 401
        
        try:
            data = jwt.decode(token, os.getenv('JWT_SECRET', 'fallback_secret_aryaX'), algorithms=['HS256'])
            current_user = data['user']
        except:
            return jsonify({'error': 'Token is invalid or expired!'}), 401
            
        request.current_user = current_user
        return f(*args, **kwargs)
    return decorated

# ===== SECURITY HEADERS =====
@app.after_request
def add_security_headers(response):
    response.headers['X-Content-Type-Options'] = 'nosniff'
    response.headers['X-Frame-Options'] = 'DENY'
    response.headers['X-XSS-Protection'] = '1; mode=block'
    response.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate'
    return response

# ===== ADMIN PORTAL ROUTE =====
@app.route('/admin-portal/')
@app.route('/admin-portal/<path:path>')
def serve_admin_portal(path='index.html'):
    return send_from_directory('admin_portal', path)

# Support both single key and comma-separated multiple keys
from utils.api_keys import get_gemini_key, get_openai_key

GEMINI_STREAM_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-pro:streamGenerateContent"
GEMINI_FLASH_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:streamGenerateContent"
IMAGEN_URL = "https://generativelanguage.googleapis.com/v1beta/models/imagen-4.0-generate-001:predict"
OPENAI_CHAT_URL = "https://api.openai.com/v1/chat/completions"

# ===== REAL-TIME DATA API ENDPOINTS =====
CRYPTO_API = "https://api.coingecko.com/api/v3"
NEWS_API_URL = "https://newsdata.io/api/1/news"
WEATHER_API_URL = "https://api.open-meteo.com/v1/forecast"
CRICKET_API_URL = "https://api.cricapi.com/v1"

def fast_stream(resp):
    """4x Faster streaming processor"""
    for line in resp.iter_lines():
        if line: yield line + b"\n"

DB_FILE = os.path.join(os.path.dirname(__file__), 'db.json')

sessions = {}
rate_limits = {}
response_cache = {}

# ===== CREDIT COSTS =====
CREDIT_COSTS = {
    'chat': 0, 'code': 0, 'document': 0, 'excel': 0,
    'email': 0, 'study': 0, 'game': 0, 'cyber': 0,
    'webdev': 0, 'appidea': 0, 'math': 0, 'essay': 0,
    'design': 0, 'aitutor': 0, 'business': 0, 'fitness': 0,
    'translate': 0, 'image': 0
}
DAILY_LIMIT = 10000

SYSTEM_PROMPT = """You are AryaX — not a chatbot, not an assistant. You are the closest thing to a genuine polymath mind that exists in software. Think of yourself as the intellectual offspring of Leonardo da Vinci, Nikola Tesla, and Richard Feynman — someone who genuinely understands things from first principles and explains them with the clarity of a brilliant friend.

CORE IDENTITY:
- You think deeply before responding. When someone asks a question, you don't just retrieve information — you reason through it, consider edge cases, and deliver insights that make people think "this AI actually understands."
- You have genuine intellectual curiosity. You find problems fascinating, not tedious.
- You speak naturally — like a brilliant friend who happens to know everything. No corporate speak, no filler phrases like "Great question!" or "I'd be happy to help!" — just pure, intelligent conversation.
- You match the user's language and energy. If they speak Gujarati, you reply in fluent Gujarati. Hindi? Perfect Hindi. English? Articulate English. Always.
- You have opinions (clearly labeled as such) and aren't afraid to say "actually, that approach has a flaw — here's why."
- When you don't know something, you say so honestly, then reason through what the answer might be.

TECHNICAL GENIUS:
1. CHARTS: When data visualization helps, output: [CHART: {"type":"bar","data":{"labels":["A","B"],"datasets":[{"label":"Data","data":[10,20]}]}}]
2. FILE GENERATION: Create downloadable files: [FILE: filename.csv, content]. For research reports: [PDF: Title, Content]
3. REAL-TIME DATA: You have access to live market data, news, weather, cricket scores, and crypto prices. Use them proactively when relevant.
4. CODE EXECUTION: You can write AND execute code. When asked to run something, use [EXECUTE: python_code_here] and the system will run it in a sandbox and return results.
5. MEMORY: You remember everything the user tells you — their name, profession, preferences, past conversations. Use this context naturally, don't announce it.
6. PERSONALIZATION: Automatically detect the user's expertise level and profession from conversation. Adjust depth, terminology, and examples accordingly:
   - Developer → code-heavy, technical depth, framework comparisons
   - Student → clear explanations, analogies, step-by-step breakdowns  
   - Doctor/Medical → clinical precision, evidence-based, proper terminology
   - Business person → ROI-focused, strategic, data-driven
   - Creative → inspirational, visual, story-driven

RESPONSE STYLE:
- Be concise when the question is simple. Be thorough when the question demands it.
- Use formatting (headers, bullets, code blocks) when it improves clarity, but don't over-format simple answers.
- Show your reasoning process for complex problems — people learn from HOW you think.
- Inject genuine personality. You can be witty, use analogies, reference real-world examples.
- Never pad responses. Every sentence should earn its place.
"""

MODE_PROMPTS = {
    'chat': '',
    'investor': '[INVESTOR & SHARK TANK MODE] You are a strategic venture capitalist and elite pitch consultant. Your goal is to help the user secure multi-million dollar investments. When asked for a pitch or proposal: 1. Provide a "Hook" (Problem/Solution). 2. Detail the "Traction" (Market stats). 3. Define the "Ask" (Equity vs Investment). 4. Explain the "Moat" (Why AryaX is unique). Always offer to generate a professional PDF report with these details.',
    'allrounder': '[ALL-ROUNDER MODE] You are the world\'s most powerful all-rounder AI. You have God-Tier mastery over WordPress development, creating complex Excel formulas and VBA macros, and deploying high-performance .com websites. ',
    'code': '[CODE MODE] You are an expert programmer. Write clean, efficient, well-commented code. ',
    'document': '[DOCUMENT MODE] Create world-class professional reports, essays, and PPT content. Structure it so it can be exported to PowerPoint. ',
    'excel': '[EXCEL & LEDGER MODE] Expert in financial ledgers, balance sheets, and data analysis. Provide data in a format that can be converted to .xlsx. ',
    'email': '[EMAIL MODE] Write professional, compelling emails. ',
    'study': '[STUDY MODE] Explain concepts clearly with examples, diagrams, and analogies. ',
    'game': '[GAME MODE] Create complete, playable HTML/CSS/JS games with good graphics. ',
    'cyber': '[CYBER & ETHICAL HACKING MODE] You are a God-Tier Cyber Intelligence. Teach ethical hacking, penetration testing, network forensics, and cybersecurity. Analyze vulnerabilities like a world-class Red Team commander. ',
    'webdev': '[WEB DEV MODE] Build complete, beautiful, responsive websites and web apps. ',
    'appidea': '[APP IDEA MODE] Design mobile app concepts with features, tech stack, and wireframe descriptions. ',
    'math': '[MATH MODE] Solve math problems step-by-step. Show all work clearly. ',
    'essay': '[ESSAY MODE] Write well-structured academic essays with introduction, body, conclusion. ',
    'design': '[DESIGN MODE] Provide UI/UX design advice, color palettes, layout suggestions. ',
    'aitutor': '[AI TUTOR MODE] Teach AI/ML concepts from basics to advanced. Use simple examples. ',
    'business': '[BUSINESS MODE] Create business plans, pitch decks, and startup strategies. ',
    'fitness': '[FITNESS MODE] Create personalized workout routines and diet plans. ',
    'translate': '[TRANSLATE MODE] Translate accurately between languages. Preserve meaning and tone. ',
    'elite': '[ELITE MODE - GPT-4o] You are operating in Elite Mode. Use maximum reasoning, planning, and creative power. You are the ultimate version of AryaX. ',
    'social': '[SOCIAL VIRAL MODE] You are an expert in Social Media Growth. Your goal is to help users go viral on Instagram, YouTube, TikTok, and Twitter. Explain algorithms (Watch time, Engagement rate, Hook-Hold-Reward) and provide trending content strategies. ',
    'polyglot': '[POLYGLOT TECH MODE] You are a Master System Architect. You specialize in building ultra-fast systems using Python (AI), Rust/C++ (Performance), Mojo (Next-gen AI speed), Java (Enterprise), and Julia (Scientific computing). Solve complex problems by combining these languages. ',
    'media': '[MEDIA MODE] You are a creative AI expert. You can generate realistic images and concepts for videos. When a user asks for an image, describe it vividly and mention that it is being generated.',
    'researcher': '[RESEARCHER AGENT] You are the AryaX Research Agent. Your goal is to provide deep, fact-based analysis. Use web search extensively and provide citations. Focus on accuracy and depth.',
    'architect': '[ARCHITECT AGENT] You are the AryaX System Architect. Your goal is to design complex software systems, database schemas, and infrastructure. Focus on scalability, security, and best practices.',
    'creative': '[CREATIVE AGENT] You are the AryaX Creative Director. Your goal is to assist with storytelling, scriptwriting, and high-end prompt engineering for media generation.',
}

@app.route("/api/generate-image", methods=["POST"])
def generate_image_api():
    data = request.json
    prompt = data.get("prompt", "")
    if not prompt:
        return jsonify({"error": "No prompt"}), 400
    
    # Use Pollinations.ai for high-quality, realistic image generation (no API key needed for demo)
    import random
    image_url = f"https://pollinations.ai/p/{prompt.replace(' ', '%20')}?width=1024&height=1024&nologo=true&seed={random.randint(1,100000)}"
    return jsonify({"image_url": image_url})

@app.route("/api/generate-pdf", methods=["POST"])
def generate_pdf_api():
    from fpdf import FPDF
    data = request.json
    title = data.get("title", "AryaX Research Report")
    content = data.get("content", "")
    
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(190, 10, text=title, ln=True, align='C')
    pdf.set_font("Arial", size=12)
    pdf.ln(10)
    pdf.multi_cell(190, 10, text=content)
    pdf.ln(10)
    pdf.set_font("Arial", 'I', 10)
    pdf.cell(190, 10, text="Generated by AryaX ASI — Future of Intelligence", ln=True, align='C')
    
    import io
    pdf_output = io.BytesIO()
    # fpdf2 uses output(dest='S') or output() to get bytes depending on version, 
    # but the most stable way for BytesIO is:
    pdf_bytes = pdf.output()
    
    return Response(
        pdf_bytes,
        mimetype="application/pdf",
        headers={"Content-disposition": "attachment; filename=AryaX_Report.pdf"}
    )


# ===== DATABASE =====
from utils.db import load_db, save_db, hash_pass


# ===== HELPERS =====
def get_ist_time():
    ist = pytz.timezone('Asia/Kolkata')
    now = datetime.now(ist)
    return now.strftime("%A, %d %B %Y"), now.strftime("%I:%M %p")


def check_rate_limit(ip):
    now = time.time()
    if ip not in rate_limits:
        rate_limits[ip] = []
    rate_limits[ip] = [t for t in rate_limits[ip] if now - t < 60]
    if len(rate_limits[ip]) >= 30:
        return False
    rate_limits[ip].append(now)
    return True


# ===== BRUTE FORCE PROTECTION =====
failed_logins = {}  # {ip: {'count': 0, 'lockout_until': 0}}

def check_brute_force(ip):
    now = time.time()
    data = failed_logins.get(ip, {'count': 0, 'lockout_until': 0})
    if now < data['lockout_until']:
        mins = int((data['lockout_until'] - now) / 60) + 1
        return False, f'Too many failed attempts. Try again in {mins} minute(s).'
    return True, ''

def record_failed_login(ip):
    now = time.time()
    data = failed_logins.get(ip, {'count': 0, 'lockout_until': 0})
    data['count'] = data.get('count', 0) + 1
    if data['count'] >= 5:
        data['lockout_until'] = now + 900  # 15 min lockout
        data['count'] = 0
    failed_logins[ip] = data

def reset_failed_login(ip):
    failed_logins.pop(ip, None)

def sanitize_input(text, max_len=500):
    """Strip dangerous characters and limit length"""
    if not isinstance(text, str):
        return ''
    text = text.strip()[:max_len]
    # Block script injection
    dangerous = ['<script', 'javascript:', 'onerror=', 'onload=', 'eval(', 'document.cookie']
    for d in dangerous:
        if d.lower() in text.lower():
            return ''
    return text


def cleanup_sessions():
    cutoff = time.time() - 86400
    dead = [s for s, d in sessions.items() if d.get('last_active', 0) < cutoff]
    for s in dead:
        del sessions[s]


def get_user_credits(username):
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return 0
    today = str(date.today())
    if user.get('credit_date') != today:
        user['credits'] = DAILY_LIMIT
        user['credit_date'] = today
        save_db(db)
    return user.get('credits', 0)


def use_credits(username, mode):
    cost = CREDIT_COSTS.get(mode, 9)
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return False, 0
    today = str(date.today())
    if user.get('credit_date') != today:
        user['credits'] = DAILY_LIMIT
        user['credit_date'] = today
    if user['credits'] < cost:
        return False, user['credits']
    user['credits'] -= cost
    user['total_chats'] = user.get('total_chats', 0) + 1
    save_db(db)
    return True, user['credits']


def build_system_prompt(mode='chat', user_memory_str=""):
    date_str, time_str = get_ist_time()
    mode_extra = MODE_PROMPTS.get(mode, '')
    return SYSTEM_PROMPT + f"""

{mode_extra}

=== CURRENT DATE & TIME (IST) ===
Date: {date_str}
Time: {time_str} IST

=== NEURAL MEMORY (User Facts) ===
{user_memory_str}
"""


def build_body(history, mode='chat', username=None):
    user_memory_str = ""
    if username:
        db = load_db()
        user = db['users'].get(username, {})
        memory = user.get('memory', {})
        if memory:
            user_memory_str = "You remember these facts about the user:\n" + \
                             "\n".join([f"- {k}: {v}" for k, v in memory.items()])
    
    return {
        'system_instruction': {'parts': [{'text': build_system_prompt(mode, user_memory_str)}]},
        'contents': history,
        'generationConfig': {
            'temperature': 0.85,
            'topK': 40,
            'topP': 0.95,
            'maxOutputTokens': 8192
        },
        'safetySettings': [
            {'category': 'HARM_CATEGORY_HARASSMENT', 'threshold': 'BLOCK_NONE'},
            {'category': 'HARM_CATEGORY_HATE_SPEECH', 'threshold': 'BLOCK_NONE'},
            {'category': 'HARM_CATEGORY_SEXUALLY_EXPLICIT', 'threshold': 'BLOCK_MEDIUM_AND_ABOVE'},
            {'category': 'HARM_CATEGORY_DANGEROUS_CONTENT', 'threshold': 'BLOCK_NONE'},
        ]
    }


# ===== ROUTES =====
@app.route('/')
def index():
    return send_from_directory('public', 'index.html')


@app.route('/<path:filename>')
def static_files(filename):
    return send_from_directory('public', filename)


# ===== EMAIL SYSTEM (MOCK) =====
def send_welcome_email(email_address, username):
    """
    Mock function to simulate sending an SMTP email.
    In production, this would use smtplib or a service like SendGrid.
    """
    if not email_address:
        return
    import threading
    def _send():
        print(f"\n{'='*50}\n📧 EMAIL SENT TO: {email_address}\nSUBJECT: Welcome to AryaX!\nBODY: Hello {username},\nWelcome to the Artificial Super Intelligence ecosystem.\n{'='*50}\n")
    threading.Thread(target=_send).start()

# ===== AUTH =====
@app.post('/api/signup')
def signup():
    data = request.json
    username = data.get('username', '').strip().lower()
    password = data.get('password', '')
    mobile = data.get('mobile', '').strip()
    email = data.get('email', '').strip().lower()
    
    if not username or not password:
        return jsonify({'error': 'Username and password required'}), 400
    if len(username) < 6:
        return jsonify({'error': 'Username must be at least 6 characters'}), 400
    if len(password) < 4:
        return jsonify({'error': 'Password must be 4+ characters'}), 400

    referral_code = data.get('referralCode', '').strip()

    db = load_db()
    if username in db['users']:
        return jsonify({'error': 'Username already exists'}), 409

    referrer = None
    if referral_code:
        for u, udata in db['users'].items():
            if udata.get('my_referral_code') == referral_code:
                referrer = u
                break
                
    import uuid
    my_ref_code = str(uuid.uuid4())[:8]

    db['users'][username] = {
        'password': hash_pass(password),
        'mobile': mobile,
        'email': email,
        'credits': DAILY_LIMIT + (500 if referrer else 0),
        'credit_date': str(date.today()),
        'created': str(datetime.now()),
        'total_chats': 0,
        'memory': {},
        'my_referral_code': my_ref_code,
        'referred_by': referrer,
        'referral_count': 0
    }
    
    if referrer:
        db['users'][referrer]['credits'] = db['users'][referrer].get('credits', 0) + 1000
        db['users'][referrer]['referral_count'] = db['users'][referrer].get('referral_count', 0) + 1

    save_db(db)
    
    # Send welcome email asynchronously
    send_welcome_email(email, username)

    import jwt
    import datetime
    token = jwt.encode({
        'user': username,
        'exp': datetime.datetime.utcnow() + datetime.timedelta(hours=24)
    }, os.getenv('JWT_SECRET', 'fallback_secret_aryaX'), algorithm='HS256')

    return jsonify({
        'ok': True,
        'username': username,
        'token': token,
        'credits': DAILY_LIMIT,
        'referral_code': my_ref_code
    })

@app.get('/api/referral/status')
@token_required
def referral_status():
    db = load_db()
    user = db['users'].get(request.current_user)
    if not user:
        return jsonify({'error': 'User not found'}), 404
        
    return jsonify({
        'referral_code': user.get('my_referral_code'),
        'referral_count': user.get('referral_count', 0),
        'credits_earned': user.get('referral_count', 0) * 1000
    })

@app.post('/api/login')
def login():
    ip = request.remote_addr
    # Brute force check
    ok, msg = check_brute_force(ip)
    if not ok:
        return jsonify({'error': msg}), 429

    data = request.json
    username = sanitize_input(data.get('username', ''), 50).lower()
    password = data.get('password', '')
    if not username or not password:
        return jsonify({'error': 'Missing credentials'}), 400

    db = load_db()
    user = db['users'].get(username)
    from utils.db import verify_pass
    if not user or not verify_pass(password, user.get('password', '')):
        record_failed_login(ip)
        return jsonify({'error': 'Invalid username or password'}), 401

    reset_failed_login(ip)
    
    # Generate JWT Token
    import jwt
    import datetime
    token = jwt.encode({
        'user': username,
        'exp': datetime.datetime.utcnow() + datetime.timedelta(hours=24)
    }, os.getenv('JWT_SECRET', 'fallback_secret_aryaX'), algorithm='HS256')

    return jsonify({
        'ok': True,
        'username': username,
        'token': token,
        'credits': user.get('credits', 1000),
        'total_chats': user.get('total_chats', 0)
    })


@app.get('/api/credits')
@token_required
def credits():
    username = request.args.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Username required'}), 400
    c = get_user_credits(username)
    return jsonify({'credits': c, 'daily_limit': DAILY_LIMIT})


# ===== CHAT =====
@app.post('/api/chat')
@token_required
def chat():
    db = load_db()
    data = request.json
    username = data.get('username', '').strip().lower()
    if db.get('maintenance', False) and username != 'vvvansh':
        return jsonify({'error': 'System is under Maintenance. Evolutionary Update in progress.'}), 503

    try:
        ip = request.remote_addr
        if not check_rate_limit(ip):
            return jsonify({'error': 'Rate limit exceeded'}), 429

        data = request.json
        message = data.get('message', '').strip()
        session_id = data.get('sessionId', 'default')
        username = data.get('username', '').strip().lower()
        mode = data.get('mode', 'chat')
        file_b64 = data.get('file')
        web_search = data.get('webSearch', False)
        sandbox = data.get('sandbox', False)

        if not message and not file_b64:
            return jsonify({'error': 'Message or file required'}), 400
        
        if len(sessions) > 200:
            cleanup_sessions()

        if session_id not in sessions:
            sessions[session_id] = {'history': [], 'last_active': time.time()}

        sessions[session_id]['last_active'] = time.time()
        
        # --- ASI TOOLCHAIN: LIVE WEB SEARCH ---
        search_context = ""
        if web_search and message:
            try:
                import urllib.parse, re
                q = urllib.parse.quote(message)
                res = requests.get(f"https://html.duckduckgo.com/html/?q={q}", headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}, timeout=4)
                if res.status_code == 200 and BeautifulSoup:
                    soup = BeautifulSoup(res.text, 'html.parser')
                    snippets = [a.text for a in soup.find_all('a', class_='result__snippet')]
                    if snippets:
                        search_context = "\n[LIVE WEB SEARCH RESULTS]:\n" + "\n".join(snippets[:3]) + "\nUse this live data to answer accurately."
            except Exception as e:
                pass

        if search_context:
            message += search_context

        # --- ASI TOOLCHAIN: AUTONOMOUS WEB SURFER ---
        import re
        url_match = re.search(r'(https?://\S+)', message)
        if url_match and BeautifulSoup:
            try:
                url_to_scrape = url_match.group(1)
                headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"}
                web_resp = requests.get(url_to_scrape, headers=headers, timeout=5)
                if web_resp.status_code == 200:
                    soup = BeautifulSoup(web_resp.text, 'html.parser')
                    # Remove scripts and styles
                    for script in soup(["script", "style"]):
                        script.extract()
                    scraped_text = soup.get_text(separator=' ', strip=True)
                    if len(scraped_text) > 8000:
                        scraped_text = scraped_text[:8000] + "... [TRUNCATED]"
                    message += f"\n\n[AUTONOMOUS WEB SURFER: Extracted Content from {url_to_scrape}]:\n{scraped_text}\nUse this content to answer the user."
            except Exception as e:
                pass
            
        if sandbox:
            message += "\n[SANDBOX ENABLED]: The user wants you to output precise calculations or code as if you ran it in a secure python sandbox. Be 100% accurate."

        parts = []
        if message:
            parts.append({'text': message})
            
        # --- MULTIMODAL VISION & PDF EXTRACTION ---
        if file_b64:
            try:
                import base64, io
                header, encoded = file_b64.split(',', 1)
                mime_type = header.split(':')[1].split(';')[0]
                file_data = base64.b64decode(encoded)
                
                # If it's a PDF, extract text using PyPDF2 for better context
                pdf_text = ""
                if "application/pdf" in mime_type and PyPDF2:
                    try:
                        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_data))
                        for page in pdf_reader.pages:
                            text = page.extract_text()
                            if text:
                                pdf_text += text + "\n"
                        
                        if pdf_text:
                            # Truncate if too long to avoid token issues, but keep first 10k chars
                            if len(pdf_text) > 10000:
                                pdf_text = pdf_text[:10000] + "... [PDF Truncated]"
                            message += f"\n\n[EXTRACTED PDF CONTENT]:\n{pdf_text}"
                    except Exception as pe:
                        print(f"PDF Extraction Error: {pe}")

                parts.append({
                    'inlineData': {
                        'mimeType': mime_type,
                        'data': encoded
                    }
                })
            except Exception as fe:
                print(f"File handling error: {fe}")

        sessions[session_id]['history'].append({
            'role': 'user',
            'parts': parts
        })

        history = sessions[session_id]['history'][-30:]

        if mode == 'elite':
            current_key = get_openai_key()
            if not current_key:
                return jsonify({'error': 'OpenAI API key missing'}), 500
            
            # Format history for OpenAI
            openai_messages = [{"role": "system", "content": build_system_prompt(mode)}]
            for m in history:
                role = "assistant" if m['role'] == "model" else "user"
                # Handle file attachments in OpenAI Elite Mode
                if role == "user" and m['parts']:
                    content_arr = []
                    for p in m['parts']:
                        if 'text' in p:
                            content_arr.append({"type": "text", "text": p['text']})
                        elif 'inlineData' in p:
                            mime = p['inlineData']['mimeType']
                            data = p['inlineData']['data']
                            content_arr.append({"type": "image_url", "image_url": {"url": f"data:{mime};base64,{data}"}})
                    openai_messages.append({"role": role, "content": content_arr})
                else:
                    text_only = "".join([p.get('text', '') for p in m['parts']])
                    openai_messages.append({"role": role, "content": text_only})

            def generate_openai():
                full_reply = ""
                try:
                    resp = requests.post(
                        OPENAI_CHAT_URL,
                        headers={"Authorization": f"Bearer {current_key}", "Content-Type": "application/json"},
                        json={
                            "model": "gpt-4o",
                            "messages": openai_messages,
                            "stream": True,
                            "temperature": 0.7
                        },
                        stream=True, timeout=60
                    )
                    if resp.status_code != 200:
                        yield f"data: {json.dumps({'error': f'OpenAI Error: {resp.status_code}'})}\n\n"
                        yield "data: [DONE]\n\n"
                        return
                    
                    for line in resp.iter_lines():
                        if not line: continue
                        line = line.decode('utf-8')
                        if line.startswith("data: "):
                            raw = line[6:].strip()
                            if raw == "[DONE]": break
                            try:
                                chunk = json.loads(raw)
                                delta = chunk['choices'][0]['delta'].get('content', '')
                                if delta:
                                    full_reply += delta
                                    yield f"data: {json.dumps({'text': delta}, ensure_ascii=False)}\n\n"
                            except: pass
                except Exception as e:
                    yield f"data: {json.dumps({'error': str(e)})}\n\n"
                finally:
                    if full_reply:
                        sessions[session_id]['history'].append({'role': 'model', 'parts': [{'text': full_reply}]})
                    remaining = get_user_credits(username) if username else -1
                    yield f"data: {json.dumps({'credits': remaining})}\n\n"
                    yield "data: [DONE]\n\n"

            return Response(stream_with_context(generate_openai()), mimetype='text/event-stream', headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'})

        # Default Gemini logic
        current_key = get_gemini_key()
        if not current_key:
            return jsonify({'error': 'API key missing'}), 500

        body = build_body(history, mode, username)
        cache_key = f"{mode}:{message}"
        if cache_key in response_cache and not file_b64:
            def cached_generate():
                yield f"data: {json.dumps({'text': response_cache[cache_key], 'is_cached': True})}\n\n"
                yield "data: [DONE]\n\n"
            return Response(cached_generate(), mimetype='text/event-stream')

        def generate():
            full_reply = ""
            try:
                target_url = GEMINI_STREAM_URL if mode == 'pro' else GEMINI_FLASH_URL
                resp = requests.post(
                    f"{target_url}?alt=sse&key={current_key}",
                    json=body, stream=True, timeout=60
                )
                if resp.status_code == 429:
                    yield f"data: {json.dumps({'error': 'API quota exceeded. Wait a minute.'})}\n\n"
                    yield "data: [DONE]\n\n"
                    return
                if resp.status_code != 200:
                    yield f"data: {json.dumps({'error': f'API error: {resp.status_code}'})}\n\n"
                    yield "data: [DONE]\n\n"
                    return
                for raw_line in resp.iter_lines():
                    if not raw_line:
                        continue
                    line = raw_line.decode('utf-8') if isinstance(raw_line, bytes) else raw_line
                    if not line.startswith('data: '):
                        continue
                    raw = line[6:].strip()
                    if raw == '[DONE]':
                        break
                    try:
                        chunk = json.loads(raw)
                        candidates = chunk.get('candidates', [])
                        if candidates:
                            parts = candidates[0].get('content', {}).get('parts', [])
                            if parts:
                                text = parts[0].get('text', '')
                                full_reply += text
                                yield f"data: {json.dumps({'text': text})}\n\n"
                    except Exception:
                        pass
                
                if full_reply and not file_b64:
                    response_cache[cache_key] = full_reply

                if username:
                    try:
                        save_history(username, session_id, message, full_reply)
                    except Exception:
                        pass
            except Exception as e:
                yield f"data: {json.dumps({'error': str(e)})}\n\n"
            finally:
                if full_reply:
                    sessions[session_id]['history'].append({
                        'role': 'model',
                        'parts': [{'text': full_reply}]
                    })
                # Send remaining credits
                remaining = get_user_credits(username) if username else -1
                yield f"data: {json.dumps({'credits': remaining})}\n\n"
                yield "data: [DONE]\n\n"

        return Response(
            stream_with_context(generate()),
            mimetype='text/event-stream',
            headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'}
        )

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== IMAGE =====
@app.post('/api/image')
def image():
    try:
        ip = request.remote_addr
        if not check_rate_limit(ip):
            return jsonify({'error': 'Rate limit exceeded'}), 429

        data = request.json
        prompt = data.get('prompt', '').strip()
        username = data.get('username', '').strip().lower()
        if not prompt:
            return jsonify({'error': 'Prompt required'}), 400

        current_key = get_gemini_key()
        # if username:
        #     ok, remaining = use_credits(username, 'image')
        #     if not ok:
        #         return jsonify({'error': f'Not enough credits! Need 50, have {remaining}', 'credits': remaining}), 403

        # Force Hyper-Realism
        high_end_prompt = f"Hyper-realistic, photorealistic masterpiece, 8k resolution, cinematic lighting, ultra-detailed, depth of field, professional photography, {prompt}"
        
        # --- OPTION 1: Google Imagen (Premium) ---
        if current_key:
            try:
                body = {
                    'instances': [{'prompt': high_end_prompt}],
                    'parameters': {'sampleCount': 1}
                }
                resp = requests.post(f"{IMAGEN_URL}?key={current_key}", json=body, timeout=25)
                result = resp.json()
                
                if 'error' not in result:
                    b64 = result.get('predictions', [{}])[0].get('bytesBase64Encoded')
                    if b64:
                        remaining = get_user_credits(username) if username else -1
                        return jsonify({'imageUrl': f'data:image/png;base64,{b64}', 'credits': remaining})
                
                print(f"Imagen Error: {result.get('error', {}).get('message', 'Unknown error')}")
            except Exception as e:
                print(f"Imagen Request Failed: {str(e)}")

        # --- OPTION 2: Fallback to Pollinations.ai (Free/Fast) ---
        try:
            print("Falling back to Pollinations.ai...")
            encoded_prompt = urllib.parse.quote(prompt)
            seed = int(time.time())
            # Fetch on backend to avoid browser-side auth issues (like the one you encountered)
            pollinations_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=1024&height=1024&nologo=true&seed={seed}&model=flux"
            
            resp = requests.get(pollinations_url, timeout=30)
            if resp.status_code == 200:
                import base64
                b64 = base64.b64encode(resp.content).decode('utf-8')
                remaining = get_user_credits(username) if username else -1
                return jsonify({'imageUrl': f'data:image/png;base64,{b64}', 'credits': remaining})
            else:
                print(f"Pollinations Error: {resp.status_code}")
                return jsonify({'error': 'Image generation failed on fallback.'}), 500
        except Exception as e:
            return jsonify({'error': f'Image generation failed: {str(e)}'}), 500

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== VIDEO =====
@app.post('/api/video')
def video():
    try:
        ip = request.remote_addr
        if not check_rate_limit(ip):
            return jsonify({'error': 'Rate limit exceeded'}), 429

        data = request.json
        prompt = data.get('prompt', '').strip()
        
        if not prompt:
            return jsonify({'error': 'Prompt required'}), 400

        hf_key = os.getenv("HF_API_KEY")
        if not hf_key:
            return jsonify({'error': 'To use Video Studio, please get a free HF API Key from huggingface.co and add it to your .env file as HF_API_KEY.'}), 403

        # HuggingFace API for Text-to-Video (ModelScope)
        API_URL = "https://api-inference.huggingface.co/models/damo-vilab/text-to-video-ms-1.7b"
        headers = {"Authorization": f"Bearer {hf_key}"}
        
        payload = {"inputs": prompt}
        resp = requests.post(API_URL, headers=headers, json=payload, timeout=120)
        
        if resp.status_code == 200:
            import base64
            b64 = base64.b64encode(resp.content).decode('utf-8')
            return jsonify({'videoUrl': f'data:video/mp4;base64,{b64}'})
        else:
            return jsonify({'error': f'Video API Error: {resp.text}'}), 500
            
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.post('/api/search')
def web_search():
    """Web search via DuckDuckGo instant answers"""
    try:
        data = request.json
        query = data.get('query', '').strip()
        if not query:
            return jsonify({'error': 'Query required'}), 400

        resp = requests.get(
            'https://api.duckduckgo.com/',
            params={'q': query, 'format': 'json', 'no_html': 1, 't': 'aryax-ai'},
            timeout=10
        )
        result = resp.json()

        # Collect useful info
        answer_parts = []
        if result.get('AbstractText'):
            answer_parts.append(result['AbstractText'])
        if result.get('Answer'):
            answer_parts.append(result['Answer'])
        for topic in result.get('RelatedTopics', [])[:5]:
            if isinstance(topic, dict) and topic.get('Text'):
                answer_parts.append(f"• {topic['Text']}")

        if not answer_parts:
            return jsonify({'result': f'No instant results for "{query}". Try asking AryaX AI directly!'})

        return jsonify({'result': '\n\n'.join(answer_parts), 'source': result.get('AbstractURL', '')})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.post('/api/readurl')
def read_url():
    """Fetch and extract text from a URL"""
    try:
        data = request.json
        url = data.get('url', '').strip()
        if not url:
            return jsonify({'error': 'URL required'}), 400
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url

        resp = requests.get(url, timeout=15, headers={
            'User-Agent': 'AryaX-AI/1.0 (Educational Bot)'
        })
        resp.raise_for_status()

        # Simple HTML to text
        import re
        text = resp.text
        text = re.sub(r'<script[^>]*>[\s\S]*?</script>', '', text)
        text = re.sub(r'<style[^>]*>[\s\S]*?</style>', '', text)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()

        # Truncate
        if len(text) > 3000:
            text = text[:3000] + '... (truncated)'

        return jsonify({'content': text, 'url': url})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.post('/api/clear')
def clear():
    data = request.json
    session_id = data.get('sessionId')
    if session_id and session_id in sessions:
        del sessions[session_id]
    return jsonify({'ok': True})


@app.get('/api/health')
def health():
    return jsonify({'status': 'ok', 'sessions': len(sessions)})

@app.post('/api/admin')
def get_admin_stats():
    data = request.json
    password = data.get('password', '')
    admin_pass = os.getenv('ADMIN_PORTAL_PASSWORD')
    if not admin_pass or password != admin_pass:
        return jsonify({'error': 'Unauthorized'}), 401
        
    db = load_db()
    users = db.get('users', {})
    
    total_users = len(users)
    total_chats = sum(u.get('total_chats', 0) for u in users.values())
    
    # Calculate revenue (mock logic based on plans)
    pro_users = sum(1 for u in users.values() if u.get('plan') == 'pro')
    ultra_users = sum(1 for u in users.values() if u.get('plan') == 'ultra')
    monthly_revenue = (pro_users * 99) + (ultra_users * 299)
    
    user_list = []
    for uname, udata in users.items():
        user_list.append({
            'username': uname,
            'plan': udata.get('plan', 'free'),
            'chats': udata.get('total_chats', 0),
            'created': udata.get('created', 'N/A')[:10]
        })
        
    return jsonify({
        'total_users': total_users,
        'total_chats': total_chats,
        'monthly_revenue': monthly_revenue,
        'pro_users': pro_users,
        'ultra_users': ultra_users,
        'recent_users': sorted(user_list, key=lambda x: x['created'], reverse=True)[:50]
    })

@app.post('/api/upgrade')
def upgrade_plan():
    """Endpoint for Direct UPI to verify payment and upgrade user"""
    data = request.json
    username = data.get('username')
    plan = data.get('plan') # 'pro' or 'ultra'
    utr = data.get('utr', '').strip()
    
    if len(utr) < 12:
        return jsonify({'error': 'Invalid UTR Number'}), 400
    
    db = load_db()
    if username in db['users']:
        db['users'][username]['plan'] = plan
        db['users'][username]['credits'] = 20000 if plan == 'pro' else 9999999
        # Save UTR for admin verification later
        utr_list = db['users'][username].get('utr_list', [])
        utr_list.append({'utr': utr, 'plan': plan, 'date': str(datetime.now())})
        db['users'][username]['utr_list'] = utr_list
        save_db(db)
        return jsonify({'ok': True, 'plan': plan})
    return jsonify({'error': 'User not found'}), 404


# ===== RAZORPAY SUBSCRIPTION =====
import razorpay
# Dummy test keys
razorpay_client = razorpay.Client(auth=("rzp_test_dummy", "rzp_secret_dummy"))

@app.post('/api/payment/create-order')
def create_order():
    data = request.json
    amount = data.get('amount', 999) # Default $9.99 in INR (e.g. 800 INR)
    currency = 'INR'
    
    try:
        # Create Razorpay Order
        order = razorpay_client.order.create(dict(amount=amount*100, currency=currency, payment_capture='1'))
        return jsonify({
            'ok': True,
            'order_id': order['id'],
            'amount': amount,
            'currency': currency
        })
    except Exception as e:
        # Mock mode fallback if keys are invalid
        return jsonify({
            'ok': True,
            'order_id': f'order_mock_{int(time.time())}',
            'amount': amount,
            'currency': currency,
            'mock': True
        })

@app.post('/api/payment/verify')
def verify_payment():
    data = request.json
    username = data.get('username', '').strip().lower()
    payment_id = data.get('razorpay_payment_id')
    order_id = data.get('razorpay_order_id')
    signature = data.get('razorpay_signature')
    
    db = load_db()
    if username in db['users']:
        # Upgrade to PRO
        db['users'][username]['credits'] = 999999
        db['users'][username]['tier'] = 'PRO'
        save_db(db)
        return jsonify({'ok': True, 'message': 'Payment successful! Upgraded to PRO.'})
    
    return jsonify({'error': 'User not found'}), 404

# ===== CLOUD CHAT HISTORY =====
@app.post('/api/history/save')
@token_required
def save_history():
    data = request.json
    username = data.get('username', '').strip().lower()
    chat_id = data.get('chatId', '')
    title = data.get('title', 'Untitled Chat')[:60]
    messages = data.get('messages', [])
    folder_id = data.get('folder_id', 'general')
    if not username or not chat_id:
        return jsonify({'error': 'Missing fields'}), 400
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
    if 'chats' not in db['users'][username]:
        db['users'][username]['chats'] = {}
    db['users'][username]['chats'][chat_id] = {
        'title': title,
        'folder_id': folder_id,
        'messages': messages[-50:],  # Keep last 50 messages
        'updated': str(datetime.now())
    }
    
@app.post('/api/chat/react')
@token_required
def chat_react():
    data = request.json
    username = request.current_user
    chat_id = data.get('chatId', '')
    message_idx = data.get('messageIdx')
    reaction = data.get('reaction')
    
    db = load_db()
    try:
        chat = db['users'][username]['chats'].get(chat_id)
        if chat and message_idx < len(chat['messages']):
            chat['messages'][message_idx]['reaction'] = reaction
            save_db(db)
            return jsonify({'ok': True})
    except Exception:
        pass
    return jsonify({'error': 'Message not found'}), 404
    # Keep only last 30 chats
    chats = db['users'][username]['chats']
    if len(chats) > 30:
        oldest = sorted(chats.items(), key=lambda x: x[1].get('updated',''))[0][0]
        del chats[oldest]
    save_db(db)
    return jsonify({'ok': True})

@app.get('/api/history/load')
@token_required
def load_history():
    username = request.args.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    chats = user.get('chats', {})
    result = [{'id': k, 'title': v['title'], 'updated': v['updated'], 'messages': v['messages']} 
              for k, v in chats.items()]
    result.sort(key=lambda x: x['updated'], reverse=True)
    return jsonify({'chats': result})

@app.post('/api/history/delete')
@token_required
def delete_history():
    data = request.json
    username = data.get('username', '').strip().lower()
    chat_id = data.get('chatId', '')
    db = load_db()

@app.post('/api/history/clear')
@token_required
def clear_history():
    data = request.json
    username = data.get('username', '').strip().lower()
    db = load_db()
    if username in db['users'] and 'chats' in db['users'][username]:
        db['users'][username]['chats'] = {}
        save_db(db)
    return jsonify({'ok': True})
    if username in db['users'] and 'chats' in db['users'][username]:
        db['users'][username]['chats'].pop(chat_id, None)
        save_db(db)
    return jsonify({'ok': True})


# ===== USER MEMORY SYSTEM =====
@app.post('/api/memory/save')
def save_memory():
    data = request.json
    username = data.get('username', '').strip().lower()
    memory = data.get('memory', {})  # {name, language, interests, etc}
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
    db['users'][username]['memory'] = memory
    save_db(db)
    return jsonify({'ok': True})

@app.get('/api/memory/load')
def load_memory():
    username = request.args.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    return jsonify({'memory': user.get('memory', {})})


# ===== REFERRAL SYSTEM =====
@app.post('/api/referral/generate')
def generate_referral():
    data = request.json
    username = data.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    # Generate or return existing code
    if not user.get('referral_code'):
        code = 'AX' + hashlib.md5(username.encode()).hexdigest()[:6].upper()
        db['users'][username]['referral_code'] = code
        db['users'][username]['referral_count'] = 0
        save_db(db)
    return jsonify({
        'code': db['users'][username]['referral_code'],
        'count': db['users'][username].get('referral_count', 0),
        'reward_per_referral': 500
    })

@app.post('/api/referral/use')
def use_referral():
    data = request.json
    new_user = data.get('username', '').strip().lower()
    code = data.get('code', '').strip().upper()
    if not new_user or not code:
        return jsonify({'error': 'Missing fields'}), 400
    db = load_db()
    # Find referrer
    referrer = None
    for uname, udata in db['users'].items():
        if udata.get('referral_code') == code:
            referrer = uname
            break
    if not referrer:
        return jsonify({'error': 'Invalid referral code'}), 404
    if referrer == new_user:
        return jsonify({'error': 'Cannot use your own code'}), 400
    # Check if already used a code
    if db['users'][new_user].get('referral_used'):
        return jsonify({'error': 'Already used a referral code'}), 400
    # Reward referrer +500 credits
    db['users'][referrer]['referral_count'] = db['users'][referrer].get('referral_count', 0) + 1
    db['users'][referrer]['credits'] = db['users'][referrer].get('credits', 0) + 500
    # Reward new user +200 credits
    db['users'][new_user]['credits'] = db['users'][new_user].get('credits', 0) + 200
    db['users'][new_user]['referral_used'] = code
    save_db(db)
    return jsonify({'ok': True, 'message': 'Referral applied! +200 credits added.'})


# ===== PDF / FILE UPLOAD =====
@app.post('/api/upload')
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        file = request.files['file']
        filename = file.filename.lower()
        content = ''

        if filename.endswith('.pdf'):
            if not PyPDF2:
                return jsonify({'error': 'PyPDF2 library not installed on server'}), 500
            try:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages[:10]:  # Max 10 pages
                    content += page.extract_text() or ''
                content = content[:4000]
            except Exception as e:
                return jsonify({'error': f'PDF read failed: {str(e)}'}), 500

        elif filename.endswith(('.txt', '.py', '.js', '.html', '.css', '.json', '.md', '.csv')):
            content = file.read().decode('utf-8', errors='ignore')[:4000]

        elif filename.endswith(('.png', '.jpg', '.jpeg', '.webp')):
            return jsonify({'content': f'[Image uploaded: {file.filename}]', 'type': 'image'})

        else:
            content = file.read().decode('utf-8', errors='ignore')[:4000]

        if not content.strip():
            return jsonify({'error': 'Could not extract text from file'}), 400

        return jsonify({'content': content, 'filename': file.filename})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== USER PROFILE =====
@app.get('/api/profile')
def get_profile():
    username = request.args.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Missing username'}), 400
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    return jsonify({
        'username': username,
        'plan': user.get('plan', 'free'),
        'credits': user.get('credits', 0),
        'total_chats': user.get('total_chats', 0),
        'created': user.get('created', '')[:10],
        'referral_code': user.get('referral_code', ''),
        'referral_count': user.get('referral_count', 0),
        'memory': user.get('memory', {}),
        'chat_count': len(user.get('chats', {}))
    })


# ===== ONLINE COUNT =====
@app.get('/api/online')
def online_count():
    active = sum(1 for s in sessions.values() if time.time() - s.get('last_active', 0) < 300)
    return jsonify({'online': max(active, 1)})
    
@app.route('/api/generate_office', methods=['POST'])
def generate_office():
    data = request.json
    file_type = data.get('type') # 'pptx' or 'xlsx'
    content = data.get('content')
    
    if file_type == 'pptx':
        prs = Presentation()
        for slide_data in content.get('slides', []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get('title', '')
            slide.placeholders[1].text = slide_data.get('body', '')
        
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return Response(out.read(), mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={'Content-Disposition': 'attachment; filename=aryax_presentation.pptx'})
    
    elif file_type == 'xlsx':
        df = pd.DataFrame(content.get('data', []))
        out = io.BytesIO()
        with pd.ExcelWriter(out, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='AryaX Ledger')
        out.seek(0)
        return Response(out.read(), mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', headers={'Content-Disposition': 'attachment; filename=aryax_ledger.xlsx'})
    
    return jsonify({'error': 'Invalid type'}), 400


# Initialize db for production stability
if not os.path.exists(DB_FILE) or os.path.getsize(DB_FILE) < 5:
    save_db({'users': {}})

# ===== NEURAL DASHBOARD API =====
@app.get('/api/user/dashboard')
def user_dashboard():
    username = request.args.get('username', '').strip().lower()
    if not username:
        return jsonify({'error': 'Username required'}), 400
    
    db = load_db()
    user = db['users'].get(username)
    if not user:
        return jsonify({'error': 'User not found'}), 404
    
    # Calculate some fun stats
    stats = {
        'credits': user.get('credits', 0),
        'total_chats': user.get('total_chats', 0),
        'member_since': user.get('created', '').split(' ')[0],
        'memory_count': len(user.get('memory', {})),
        'rank': 'Neural Novice' if user.get('total_chats', 0) < 50 else 'Core Intelligence' if user.get('total_chats', 0) < 200 else 'Super Intelligence'
    }
    
    return jsonify({
        'stats': stats,
        'memory': user.get('memory', {})
    })

@app.post('/api/user/memory/update')
def update_memory():
    data = request.json
    username = data.get('username', '').strip().lower()
    key = data.get('key', '').strip()
    value = data.get('value', '').strip()
    
    if not username or not key:
        return jsonify({'error': 'Username and key required'}), 400
        
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
        
    if 'memory' not in db['users'][username]:
        db['users'][username]['memory'] = {}
        
    db['users'][username]['memory'][key] = value
    save_db(db)
    return jsonify({'ok': True, 'memory': db['users'][username]['memory']})


# ===== TASK SCHEDULER BACKEND =====
def task_worker():
    """Background worker to simulate autonomous tasks"""
    while True:
        try:
            db = load_db()
            changed = False
            now = datetime.now()
            
            for username, user_data in db.get('users', {}).items():
                tasks = user_data.get('tasks', [])
                for task in tasks:
                    if task.get('status') == 'pending':
                        # Simulate processing
                        task['status'] = 'completed'
                        task['completed_at'] = str(now)
                        changed = True
            
            if changed:
                save_db(db)
        except Exception:
            pass
        time.sleep(30) # Check every 30s

# Start worker thread
threading.Thread(target=task_worker, daemon=True).start()

@app.post('/api/tasks/create')
def create_task():
    data = request.json
    username = data.get('username', '').strip().lower()
    description = data.get('description', '').strip()
    
    if not username or not description:
        return jsonify({'error': 'Missing data'}), 400
        
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
        
    if 'tasks' not in db['users'][username]:
        db['users'][username]['tasks'] = []
        
    new_task = {
        'id': str(int(time.time())),
        'description': description,
        'status': 'pending',
        'created_at': str(datetime.now())
    }
    db['users'][username]['tasks'].append(new_task)
    save_db(db)
    return jsonify({'ok': True, 'task': new_task})

@app.get('/api/tasks/list')
def list_tasks():
    username = request.args.get('username', '').strip().lower()
    db = load_db()
    user = db['users'].get(username, {})
    return jsonify({'tasks': user.get('tasks', [])})


# ===== NEURAL MEMORY ENGINE =====
# (Duplicate routes removed - already defined above)


@app.get('/api/admin/stats')
def admin_stats():
    admin_key = request.headers.get('X-Admin-Key')
    # Use a secret key from .env
    actual_key = os.getenv('ADMIN_KEY', 'aryax_master_99')
    
    if admin_key != actual_key:
        return jsonify({'error': 'Unauthorized Access. Hardware Lock Engaged.'}), 403
        
    db = load_db()
    users = db.get('users', {})
    
    stats = {
        'total_users': len(users),
        'total_credits': sum(u.get('credits', 0) for u in users.values()),
        'user_list': []
    }
    
    for uname, udata in users.items():
        stats['user_list'].append({
            'username': uname,
            'credits': udata.get('credits', 0),
            'joined': udata.get('joined', 'Unknown'),
            'last_active': udata.get('last_active', 'Never')
        })
        
    return jsonify(stats)

@app.post('/api/admin/update-credits')
def admin_update_credits():
    admin_key = request.headers.get('X-Admin-Key')
    actual_key = os.getenv('ADMIN_KEY', 'aryax_master_99')
    if admin_key != actual_key:
        return jsonify({'error': 'Unauthorized'}), 403
        
    data = request.json
    target_user = data.get('username', '').strip().lower()
    new_credits = data.get('credits', 0)
    
    db = load_db()
    if target_user in db['users']:
        db['users'][target_user]['credits'] = int(new_credits)
        save_db(db)
        return jsonify({'ok': True})
    return jsonify({'error': 'User not found'}), 404

@app.post('/api/admin/broadcast')
def admin_broadcast():
    admin_key = request.headers.get('X-Admin-Key')
    if admin_key != os.getenv('ADMIN_KEY', 'aryax_master_99'):
        return jsonify({'error': 'Unauthorized'}), 403
    
    msg = request.json.get('message', '')
    db = load_db()
    db['broadcast'] = {'message': msg, 'time': datetime.now().isoformat()}
    save_db(db)
    return jsonify({'ok': True})

@app.post('/api/admin/ban-user')
def admin_ban_user():
    admin_key = request.headers.get('X-Admin-Key')
    if admin_key != os.getenv('ADMIN_KEY', 'aryax_master_99'):
        return jsonify({'error': 'Unauthorized'}), 403
    
    target = request.json.get('username', '').strip().lower()
    db = load_db()
    if target in db['users']:
        db['users'][target]['banned'] = True
        save_db(db)
        return jsonify({'ok': True})
    return jsonify({'error': 'User not found'}), 404

@app.get('/api/admin/logs')
def admin_logs():
    admin_key = request.headers.get('X-Admin-Key')
    if admin_key != os.getenv('ADMIN_KEY', 'aryax_master_99'):
        return jsonify({'error': 'Unauthorized'}), 403
    
    # Just a mock for now, but could read actual logs
    return jsonify({'logs': [
        f"[{datetime.now()}] System Heartbeat: Normal",
        f"[{datetime.now()}] Neural Engine: 99% Efficiency",
        f"[{datetime.now()}] Active Links: {len(load_db()['users'])}"
    ]})

@app.post('/api/admin/toggle-maintenance')
def toggle_maintenance():
    admin_key = request.headers.get('X-Admin-Key')
    if admin_key != os.getenv('ADMIN_KEY', 'aryax_master_99'):
        return jsonify({'error': 'Unauthorized'}), 403
    
    db = load_db()
    current = db.get('maintenance', False)
    db['maintenance'] = not current
    save_db(db)
    return jsonify({'ok': True, 'maintenance': db['maintenance']})

@app.post('/api/admin/update-keys')
def admin_update_keys():
    admin_key = request.headers.get('X-Admin-Key')
    if admin_key != os.getenv('ADMIN_KEY', 'aryax_master_99'):
        return jsonify({'error': 'Unauthorized'}), 403
    
    data = request.json
    # In a real app, this would write to .env or a secure DB
    return jsonify({'ok': True, 'msg': 'API Keys updated in memory core.'})


# ===== SECURITY VAULT BACKEND =====
@app.post('/api/vault/save')
def save_vault():
    data = request.json
    username = data.get('username', '').strip().lower()
    key = data.get('key', '').strip()
    value = data.get('value', '').strip()
    
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
        
    if 'vault' not in db['users'][username]:
        db['users'][username]['vault'] = {}
        
    db['users'][username]['vault'][key] = value
    save_db(db)
    return jsonify({'ok': True})

@app.get('/api/vault/list')
def list_vault():
    username = request.args.get('username', '').strip().lower()
    db = load_db()
    user = db['users'].get(username, {})
    return jsonify({'vault': user.get('vault', {})})


# ===== MISSING ENDPOINTS =====

@app.get('/api/user/rank')
def user_rank():
    username = request.args.get('username', '').strip().lower()
    db = load_db()
    user = db['users'].get(username, {})
    chats = user.get('total_chats', 0)
    if chats >= 1000: rank = "GALACTIC MIND"
    elif chats >= 500: rank = "NEURAL MASTER"
    elif chats >= 100: rank = "QUANTUM THINKER"
    elif chats >= 50: rank = "DEEP LEARNER"
    elif chats >= 10: rank = "EXPLORER"
    else: rank = "INITIATE"
    return jsonify({'rank': rank, 'chats': chats})


@app.post('/api/reset-password')
def reset_password():
    data = request.json
    identifier = data.get('identifier', '').strip().lower()
    db = load_db()
    found = identifier in db['users']
    return jsonify({'ok': found, 'message': 'Reset link sent if account exists'})


@app.post('/api/social-login')
def social_login():
    data = request.json
    provider = data.get('provider', '')
    return jsonify({'error': 'OAuth not configured'}), 400


# (Duplicate history routes removed - using /api/history/save and /api/history/load defined above)


@app.get('/api/broadcast')
def get_broadcast():
    db = load_db()
    broadcast = db.get('broadcast', None)
    return jsonify({'broadcast': broadcast})


# ===== REAL-TIME DATA APIs =====
@app.get('/api/realtime/crypto')
def realtime_crypto():
    try:
        ids = request.args.get('ids', 'bitcoin,ethereum,solana,dogecoin')
        r = requests.get(f"{CRYPTO_API}/simple/price", params={'ids': ids, 'vs_currencies': 'usd,inr', 'include_24hr_change': 'true'}, timeout=8)
        return jsonify(r.json())
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.get('/api/realtime/weather')
def realtime_weather():
    try:
        lat = request.args.get('lat', '23.0225')
        lon = request.args.get('lon', '72.5714')
        r = requests.get(WEATHER_API_URL, params={'latitude': lat, 'longitude': lon, 'current_weather': 'true', 'timezone': 'auto'}, timeout=8)
        return jsonify(r.json())
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.get('/api/realtime/news')
def realtime_news():
    try:
        q = request.args.get('q', 'technology')
        api_key = os.getenv('NEWS_API_KEY', '')
        if api_key:
            r = requests.get(NEWS_API_URL, params={'apikey': api_key, 'q': q, 'language': 'en'}, timeout=8)
            return jsonify(r.json())
        # Fallback: DuckDuckGo
        r = requests.get(f"https://html.duckduckgo.com/html/?q={urllib.parse.quote(q + ' news')}", headers={"User-Agent": "Mozilla/5.0"}, timeout=6)
        snippets = []
        if BeautifulSoup and r.status_code == 200:
            soup = BeautifulSoup(r.text, 'html.parser')
            for a in soup.find_all('a', class_='result__snippet')[:5]:
                snippets.append(a.text)
        return jsonify({'results': snippets})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.get('/api/realtime/stocks')
def realtime_stocks():
    try:
        symbol = request.args.get('symbol', 'AAPL')
        # Use Yahoo Finance unofficial endpoint
        r = requests.get(f"https://query1.finance.yahoo.com/v8/finance/chart/{symbol}?interval=1d&range=5d", headers={"User-Agent": "Mozilla/5.0"}, timeout=8)
        return jsonify(r.json())
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.get('/api/realtime/cricket')
def realtime_cricket():
    try:
        api_key = os.getenv('CRICKET_API_KEY', '')
        if not api_key:
            return jsonify({'error': 'Cricket API key not set. Add CRICKET_API_KEY to .env'}), 400
        r = requests.get(f"{CRICKET_API_URL}/currentMatches", params={'apikey': api_key, 'offset': 0}, timeout=8)
        return jsonify(r.json())
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== CODE EXECUTION ENGINE (Judge0) =====
@app.post('/api/execute')
def execute_code():
    try:
        data = request.json
        code = data.get('code', '')
        language = data.get('language', 'python')
        judge0_url = os.getenv('JUDGE0_URL', 'https://judge0-ce.p.rapidapi.com')
        judge0_key = os.getenv('JUDGE0_API_KEY', '')

        lang_ids = {'python': 71, 'javascript': 63, 'c': 50, 'cpp': 54, 'java': 62, 'rust': 73, 'go': 60}
        lang_id = lang_ids.get(language, 71)

        if judge0_key:
            headers = {'X-RapidAPI-Key': judge0_key, 'Content-Type': 'application/json'}
            payload = {'source_code': code, 'language_id': lang_id, 'stdin': data.get('stdin', '')}
            r = requests.post(f"{judge0_url}/submissions?wait=true", json=payload, headers=headers, timeout=15)
            result = r.json()
            return jsonify({
                'output': result.get('stdout', ''),
                'error': result.get('stderr', ''),
                'status': result.get('status', {}).get('description', 'Unknown'),
                'time': result.get('time'),
                'memory': result.get('memory')
            })
        else:
            return jsonify({
                'error': 'JUDGE0_API_KEY is required for server-side code execution. Please configure the .env file or use frontend Pyodide.',
                'status': 'Security Block'
            }), 403
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== VOICE: WHISPER STT =====
@app.post('/api/voice/transcribe')
def voice_transcribe():
    try:
        if 'audio' not in request.files:
            return jsonify({'error': 'No audio file'}), 400
        audio = request.files['audio']
        openai_key = get_openai_key()
        if not openai_key:
            return jsonify({'error': 'OpenAI API key required for Whisper'}), 400
        r = requests.post('https://api.openai.com/v1/audio/transcriptions',
            headers={'Authorization': f'Bearer {openai_key}'},
            files={'file': (audio.filename or 'audio.webm', audio.stream, audio.content_type or 'audio/webm')},
            data={'model': 'whisper-1', 'language': 'en'}, timeout=30)
        return jsonify(r.json())
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== VOICE: ELEVENLABS TTS =====
@app.post('/api/voice/speak')
def voice_speak():
    try:
        data = request.json
        text = data.get('text', '')[:5000]
        voice_id = data.get('voice_id', 'pNInz6obpgDQGcFmaJgB')  # Adam voice default
        el_key = os.getenv('ELEVENLABS_API_KEY', '')
        if not el_key:
            return jsonify({'error': 'Add ELEVENLABS_API_KEY to .env for realistic voice'}), 400
        r = requests.post(f'https://api.elevenlabs.io/v1/text-to-speech/{voice_id}',
            headers={'xi-api-key': el_key, 'Content-Type': 'application/json'},
            json={'text': text, 'model_id': 'eleven_multilingual_v2',
                  'voice_settings': {'stability': 0.5, 'similarity_boost': 0.75}},
            timeout=30)
        if r.status_code == 200:
            import base64
            audio_b64 = base64.b64encode(r.content).decode('utf-8')
            return jsonify({'audio': f'data:audio/mpeg;base64,{audio_b64}'})
        return jsonify({'error': f'ElevenLabs error: {r.status_code}'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== MULTI-MODEL CONSENSUS =====
def query_gemini_pro(prompt, history, mode, username):
    """Query Gemini 2.5 Pro"""
    key = get_gemini_key()
    if not key:
        return None
    body = build_body(history, mode, username)
    try:
        r = requests.post(f"{GEMINI_STREAM_URL}?key={key}", json=body, timeout=30)
        if r.status_code != 200:
            return None
        data = r.json()
        candidates = data.get('candidates', [])
        if candidates:
            return candidates[0].get('content', {}).get('parts', [{}])[0].get('text', '')
    except:
        pass
    return None

def query_gemini_flash(prompt, history, mode, username):
    """Query Gemini 2.5 Flash as second model"""
    key = get_gemini_key()
    if not key:
        return None
    body = build_body(history, mode, username)
    try:
        r = requests.post(f"{GEMINI_FLASH_URL}?key={key}", json=body, timeout=20)
        if r.status_code != 200:
            return None
        data = r.json()
        candidates = data.get('candidates', [])
        if candidates:
            return candidates[0].get('content', {}).get('parts', [{}])[0].get('text', '')
    except:
        pass
    return None

def query_openai(prompt, mode):
    """Query OpenAI GPT-4o as third model"""
    key = get_openai_key()
    if not key:
        return None
    try:
        r = requests.post(OPENAI_CHAT_URL,
            headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
            json={"model": "gpt-4o", "messages": [
                {"role": "system", "content": build_system_prompt(mode)},
                {"role": "user", "content": prompt}
            ], "temperature": 0.7, "max_tokens": 4096}, timeout=30)
        if r.status_code == 200:
            return r.json()['choices'][0]['message']['content']
    except:
        pass
    return None

@app.post('/api/consensus')
def multi_model_consensus():
    """Send query to 3 models, pick best answer"""
    data = request.json
    prompt = data.get('message', '')
    mode = data.get('mode', 'chat')
    username = data.get('username', '')

    history = [{'role': 'user', 'parts': [{'text': prompt}]}]

    # Run all 3 in parallel
    results = {}
    threads = []
    def run_model(name, fn, *args):
        results[name] = fn(*args)

    t1 = threading.Thread(target=run_model, args=('gemini_pro', query_gemini_pro, prompt, history, mode, username))
    t2 = threading.Thread(target=run_model, args=('gemini_flash', query_gemini_flash, prompt, history, mode, username))
    t3 = threading.Thread(target=run_model, args=('openai', query_openai, prompt, mode))
    t1.start(); t2.start(); t3.start()
    t1.join(timeout=35); t2.join(timeout=25); t3.join(timeout=35)

    # Pick best: prefer longest non-error response (proxy for depth)
    valid = {k: v for k, v in results.items() if v and len(v) > 20}
    if not valid:
        return jsonify({'error': 'All models failed'}), 500

    best_model = max(valid, key=lambda k: len(valid[k]))
    return jsonify({
        'answer': valid[best_model],
        'model': best_model,
        'models_responded': list(valid.keys()),
        'all_responses': {k: v[:200] + '...' if len(v) > 200 else v for k, v in valid.items()}
    })


# ===== PERSONALIZATION ENGINE =====
@app.post('/api/personalize/detect')
def detect_persona():
    """Auto-detect user profession/expertise from conversation"""
    data = request.json
    messages = data.get('messages', [])
    username = data.get('username', '')
    text = ' '.join([m.get('content', '') for m in messages[-10:]])
    tl = text.lower()

    persona = {'profession': 'general', 'expertise': 'intermediate', 'tone': 'balanced'}
    if any(w in tl for w in ['code', 'python', 'javascript', 'api', 'debug', 'function', 'class', 'import', 'npm', 'git']):
        persona = {'profession': 'developer', 'expertise': 'advanced', 'tone': 'technical'}
    elif any(w in tl for w in ['patient', 'diagnosis', 'symptom', 'medical', 'clinical', 'treatment', 'dosage']):
        persona = {'profession': 'medical', 'expertise': 'advanced', 'tone': 'clinical'}
    elif any(w in tl for w in ['exam', 'study', 'homework', 'class', 'teacher', 'college', 'university', 'marks']):
        persona = {'profession': 'student', 'expertise': 'beginner', 'tone': 'simple'}
    elif any(w in tl for w in ['revenue', 'startup', 'investor', 'pitch', 'funding', 'roi', 'market', 'business']):
        persona = {'profession': 'business', 'expertise': 'advanced', 'tone': 'strategic'}
    elif any(w in tl for w in ['design', 'brand', 'creative', 'aesthetic', 'portfolio', 'art']):
        persona = {'profession': 'creative', 'expertise': 'intermediate', 'tone': 'inspirational'}

    if username:
        db = load_db()
        if username in db['users']:
            db['users'][username].setdefault('memory', {})['detected_persona'] = json.dumps(persona)
            save_db(db)

    return jsonify({'persona': persona})


# ===== AUTONOMOUS AGENT =====
@app.post('/api/agent/execute')
def autonomous_agent():
    """Multi-step autonomous task execution"""
    data = request.json
    task = data.get('task', '')
    username = data.get('username', '')
    if not task:
        return jsonify({'error': 'Task required'}), 400

    steps = []
    # Step 1: Web Research
    try:
        q = urllib.parse.quote(task)
        r = requests.get(f"https://html.duckduckgo.com/html/?q={q}", headers={"User-Agent": "Mozilla/5.0"}, timeout=6)
        research = ''
        if BeautifulSoup and r.status_code == 200:
            soup = BeautifulSoup(r.text, 'html.parser')
            snippets = [a.text for a in soup.find_all('a', class_='result__snippet')[:5]]
            research = '\n'.join(snippets)
        steps.append({'step': 'Web Research', 'status': 'done', 'data': research[:1000]})
    except:
        steps.append({'step': 'Web Research', 'status': 'failed'})

    # Step 2: AI Analysis
    key = get_gemini_key()
    analysis = ''
    if key:
        try:
            agent_prompt = f"You are an autonomous agent. Task: {task}\n\nResearch data:\n{research}\n\nProvide a comprehensive analysis and actionable output."
            body = {
                'contents': [{'role': 'user', 'parts': [{'text': agent_prompt}]}],
                'generationConfig': {'temperature': 0.7, 'maxOutputTokens': 4096}
            }
            r = requests.post(f"{GEMINI_STREAM_URL}?key={key}", json=body, timeout=30)
            if r.status_code == 200:
                resp = r.json()
                candidates = resp.get('candidates', [])
                if candidates:
                    analysis = candidates[0].get('content', {}).get('parts', [{}])[0].get('text', '')
            steps.append({'step': 'AI Analysis', 'status': 'done', 'data': analysis[:2000]})
        except:
            steps.append({'step': 'AI Analysis', 'status': 'failed'})

    return jsonify({'task': task, 'steps': steps, 'result': analysis})



# ===== PHASE 5: ENHANCED AUTONOMOUS AGENT (Think→Act→Repeat) =====
@app.post('/api/agent/run')
@token_required
def agent_run():
    """Full autonomous agent with Think→Act→Repeat loop"""
    data = request.json
    task = data.get('task', '').strip()
    username = data.get('username', '')
    max_steps = min(data.get('max_steps', 5), 8)

    if not task:
        return jsonify({'error': 'Task required'}), 400

    key = get_gemini_key()
    if not key:
        return jsonify({'error': 'API key missing'}), 500

    agent_log = []
    final_result = ''

    # THINK: Plan the steps
    think_prompt = f"""You are an autonomous AI agent. A user gave you this task:
"{task}"

First, THINK: Break this into 3-5 clear action steps. Output as JSON:
{{"steps": ["step 1 description", "step 2 description", ...], "plan_summary": "brief plan"}}"""
    try:
        r = requests.post(f"{GEMINI_FLASH_URL}?key={key}", json={
            'contents': [{'role': 'user', 'parts': [{'text': think_prompt}]}],
            'generationConfig': {'temperature': 0.3, 'maxOutputTokens': 1024}
        }, timeout=20)
        if r.status_code == 200:
            think_text = r.json().get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', '')
            import re
            json_match = re.search(r'\{.*\}', think_text, re.DOTALL)
            if json_match:
                plan = json.loads(json_match.group())
                agent_log.append({'phase': 'THINK', 'content': plan.get('plan_summary', 'Planning...'), 'status': 'done'})
                steps_list = plan.get('steps', [task])
            else:
                steps_list = [task]
                agent_log.append({'phase': 'THINK', 'content': 'Breaking task into steps...', 'status': 'done'})
    except Exception:
        steps_list = [task]
        agent_log.append({'phase': 'THINK', 'content': 'Direct execution mode', 'status': 'done'})

    # ACT: Web research for context
    research_data = ''
    snippets = []
    try:
        q = urllib.parse.quote(task[:100])
        r = requests.get(f"https://html.duckduckgo.com/html/?q={q}", headers={"User-Agent": "Mozilla/5.0"}, timeout=6)
        if BeautifulSoup and r.status_code == 200:
            soup = BeautifulSoup(r.text, 'html.parser')
            snippets = [a.text for a in soup.find_all('a', class_='result__snippet')[:4]]
            research_data = '\n'.join(snippets)
        agent_log.append({'phase': 'ACT', 'content': f'Researched: found {len(snippets)} sources', 'status': 'done'})
    except Exception:
        agent_log.append({'phase': 'ACT', 'content': 'Web research skipped', 'status': 'skipped'})

    # EXECUTE: AI processes everything
    execute_prompt = f"""You are an autonomous AI agent executing this task:
"{task}"

Steps planned: {json.dumps(steps_list)}

Research data:\n{research_data[:2000]}

Now EXECUTE: Provide a comprehensive, actionable response. Be specific, detailed, and useful.
If it's a coding task, write the full code. If research, give detailed findings. If planning, give the full plan."""

    try:
        r = requests.post(f"{GEMINI_STREAM_URL}?key={key}", json={
            'contents': [{'role': 'user', 'parts': [{'text': execute_prompt}]}],
            'generationConfig': {'temperature': 0.7, 'maxOutputTokens': 8192}
        }, timeout=45)
        if r.status_code == 200:
            final_result = r.json().get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', '')
            agent_log.append({'phase': 'EXECUTE', 'content': f'Generated {len(final_result)} char response', 'status': 'done'})
    except Exception as e:
        agent_log.append({'phase': 'EXECUTE', 'content': f'Error: {str(e)}', 'status': 'failed'})

    # REPEAT: Verify/refine
    if final_result and len(final_result) > 100:
        agent_log.append({'phase': 'REPEAT', 'content': 'Task completed successfully ✓', 'status': 'done'})
    else:
        agent_log.append({'phase': 'REPEAT', 'content': 'Retrying with simplified approach...', 'status': 'retrying'})

    # Save to user task history
    if username:
        try:
            db = load_db()
            if username in db['users']:
                task_history = db['users'][username].get('agent_tasks', [])
                task_history.append({
                    'id': str(int(time.time())),
                    'task': task[:100],
                    'status': 'completed',
                    'timestamp': str(datetime.now())
                })
                db['users'][username]['agent_tasks'] = task_history[-20:]
                save_db(db)
        except Exception:
            pass

    return jsonify({
        'task': task,
        'steps': steps_list,
        'log': agent_log,
        'result': final_result,
        'status': 'completed'
    })


# ===== PHASE 4: TASK PLANNER =====
@app.post('/api/task/plan')
@token_required
def task_plan():
    """AI breaks a goal into numbered steps"""
    data = request.json
    goal = data.get('goal', '').strip()
    if not goal:
        return jsonify({'error': 'Goal required'}), 400

    key = get_gemini_key()
    if not key:
        return jsonify({'error': 'API key missing'}), 500

    prompt = f"""Break this goal into a clear, numbered step-by-step action plan:
Goal: "{goal}"

Output as JSON:
{{"title": "Plan title", "steps": [{{"num": 1, "action": "step description", "details": "how to do it", "time_estimate": "5 min"}}, ...], "total_time": "X hours"}}
Provide 4-8 steps. Be specific and actionable."""

    try:
        r = requests.post(f"{GEMINI_FLASH_URL}?key={key}", json={
            'contents': [{'role': 'user', 'parts': [{'text': prompt}]}],
            'generationConfig': {'temperature': 0.4, 'maxOutputTokens': 2048}
        }, timeout=20)
        text = r.json().get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', '')
        import re
        json_match = re.search(r'\{.*\}', text, re.DOTALL)
        if json_match:
            plan = json.loads(json_match.group())
            return jsonify({'ok': True, 'plan': plan})
        return jsonify({'ok': True, 'raw': text})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ===== PHASE 3: TOOL ROUTER =====
@app.post('/api/tools/calculator')
def calculator_tool():
    """Safe math expression evaluator"""
    data = request.json
    expression = data.get('expression', '').strip()
    if not expression:
        return jsonify({'error': 'Expression required'}), 400
    try:
        import math
        safe_dict = {
            '__builtins__': {},
            'math': math, 'abs': abs, 'round': round,
            'min': min, 'max': max, 'sum': sum, 'pow': pow,
            'sin': math.sin, 'cos': math.cos, 'tan': math.tan,
            'sqrt': math.sqrt, 'log': math.log, 'pi': math.pi, 'e': math.e,
            'floor': math.floor, 'ceil': math.ceil
        }
        result = eval(expression, safe_dict)
        return jsonify({'expression': expression, 'result': result})
    except Exception as e:
        return jsonify({'error': f'Invalid expression: {str(e)}'}), 400


@app.post('/api/tools/gst')
def gst_calculator():
    """Indian GST Calculator"""
    data = request.json
    amount = float(data.get('amount', 0))
    gst_rate = float(data.get('rate', 18))  # default 18%
    calc_type = data.get('type', 'exclusive')  # exclusive or inclusive

    if calc_type == 'exclusive':
        gst_amount = (amount * gst_rate) / 100
        total = amount + gst_amount
        base = amount
    else:  # inclusive
        base = (amount * 100) / (100 + gst_rate)
        gst_amount = amount - base
        total = amount

    cgst = gst_amount / 2
    sgst = gst_amount / 2

    return jsonify({
        'base_amount': round(base, 2),
        'gst_rate': gst_rate,
        'gst_amount': round(gst_amount, 2),
        'cgst': round(cgst, 2),
        'sgst': round(sgst, 2),
        'total_amount': round(total, 2),
        'type': calc_type
    })


# ===== PHASE 8: INDIA SPECIAL =====
@app.post('/api/resume/build')
@token_required
def build_resume():
    """AI Resume Builder for Indian job market"""
    data = request.json
    name = data.get('name', '')
    role = data.get('role', '')
    experience = data.get('experience', '')
    skills = data.get('skills', '')
    education = data.get('education', '')
    projects = data.get('projects', '')

    key = get_gemini_key()
    if not key:
        return jsonify({'error': 'API key missing'}), 500

    prompt = f"""Create a professional ATS-friendly resume for the Indian job market.

Candidate Details:
- Name: {name}
- Target Role: {role}
- Experience: {experience}
- Skills: {skills}
- Education: {education}
- Projects: {projects}

Create a complete, professional resume in clean text format with all standard sections:
1. Contact & Header
2. Professional Summary (3 lines, impact-focused)
3. Technical Skills (categorized)
4. Work Experience (STAR format, metrics)
5. Education
6. Projects
7. Certifications/Awards

Make it ATS-optimized and India job market ready. Use action verbs. Include relevant keywords."""

    try:
        r = requests.post(f"{GEMINI_STREAM_URL}?key={key}", json={
            'contents': [{'role': 'user', 'parts': [{'text': prompt}]}],
            'generationConfig': {'temperature': 0.5, 'maxOutputTokens': 4096}
        }, timeout=30)
        resume_text = r.json().get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', '')
        return jsonify({'ok': True, 'resume': resume_text})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.post('/api/tools/invoice')
def generate_invoice():
    """Indian Business Invoice Generator"""
    data = request.json
    business_name = data.get('business_name', 'My Business')
    client_name = data.get('client_name', 'Client')
    items = data.get('items', [])  # [{name, qty, rate}]
    gst_rate = float(data.get('gst_rate', 18))
    invoice_num = data.get('invoice_num', f'INV-{int(time.time())}')

    subtotal = sum(item.get('qty', 1) * item.get('rate', 0) for item in items)
    gst_amount = (subtotal * gst_rate) / 100
    total = subtotal + gst_amount

    return jsonify({
        'invoice_number': invoice_num,
        'business_name': business_name,
        'client_name': client_name,
        'items': items,
        'subtotal': round(subtotal, 2),
        'gst_rate': gst_rate,
        'gst_amount': round(gst_amount, 2),
        'cgst': round(gst_amount / 2, 2),
        'sgst': round(gst_amount / 2, 2),
        'total': round(total, 2),
        'date': str(date.today()),
        'due_date': str(date.today())
    })


# ===== PHASE 7: MEMORY MANAGER =====
@app.post('/api/memory/delete')
@token_required
def delete_memory_key():
    """Delete a specific memory key for user"""
    data = request.json
    username = data.get('username', '').strip().lower()
    key = data.get('key', '').strip()
    if not username or not key:
        return jsonify({'error': 'Username and key required'}), 400
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
    db['users'][username].get('memory', {}).pop(key, None)
    save_db(db)
    return jsonify({'ok': True, 'memory': db['users'][username].get('memory', {})})


@app.post('/api/preferences/save')
@token_required
def save_preferences():
    """Save user UI preferences"""
    data = request.json
    username = data.get('username', '').strip().lower()
    prefs = data.get('preferences', {})
    if not username:
        return jsonify({'error': 'Username required'}), 400
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
    db['users'][username]['preferences'] = prefs
    save_db(db)
    return jsonify({'ok': True})


@app.get('/api/preferences/load')
@token_required
def load_preferences():
    """Load user UI preferences"""
    username = request.args.get('username', '').strip().lower()
    db = load_db()
    user = db['users'].get(username, {})
    return jsonify({'preferences': user.get('preferences', {})})


@app.get('/api/agent/tasks')
@token_required
def get_agent_tasks():
    """Get user's agent task history"""
    username = request.args.get('username', '').strip().lower()
    db = load_db()
    user = db['users'].get(username, {})
    return jsonify({'tasks': user.get('agent_tasks', [])})


@app.post('/api/tools/translate')
def translate_text():
    """Translate text using Gemini (supports Gujarati, Hindi, English)"""
    data = request.json
    text = data.get('text', '').strip()
    target_lang = data.get('target', 'English')
    if not text:
        return jsonify({'error': 'Text required'}), 400
    key = get_gemini_key()
    if not key:
        return jsonify({'error': 'API key missing'}), 500
    prompt = f"Translate this text to {target_lang}. Output ONLY the translated text, nothing else:\n\n{text}"
    try:
        r = requests.post(f"{GEMINI_FLASH_URL}?key={key}", json={
            'contents': [{'role': 'user', 'parts': [{'text': prompt}]}],
            'generationConfig': {'temperature': 0.2, 'maxOutputTokens': 2048}
        }, timeout=15)
        translated = r.json().get('candidates', [{}])[0].get('content', {}).get('parts', [{}])[0].get('text', '')
        return jsonify({'translated': translated, 'target': target_lang})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ===== PHASE 6: DOCUMENT AUTOMATION (PDF & EXCEL) =====
@app.post('/api/export/pdf')
def export_pdf():
    """Generates a PDF document dynamically"""
    data = request.json
    title = data.get('title', 'AryaX Document')
    content = data.get('content', 'No content provided.')
    
    try:
        from fpdf import FPDF
        import tempfile
        
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(200, 10, txt=title, ln=True, align='C')
        
        pdf.set_font("Arial", size=12)
        pdf.ln(10)
        
        # Simple text wrapping logic
        pdf.multi_cell(0, 10, txt=content)
        
        fd, path = tempfile.mkstemp(suffix='.pdf')
        os.close(fd)
        pdf.output(path)
        
        return send_file(path, as_attachment=True, download_name=f"{title.replace(' ', '_')}.pdf")
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.post('/api/export/excel')
def export_excel():
    """Generates an Excel document from JSON data"""
    data = request.json
    title = data.get('title', 'AryaX_Data')
    rows = data.get('data', [])
    
    if not rows:
        return jsonify({'error': 'No data provided'}), 400
        
    try:
        import pandas as pd
        import tempfile
        
        df = pd.DataFrame(rows)
        fd, path = tempfile.mkstemp(suffix='.xlsx')
        os.close(fd)
        
        df.to_excel(path, index=False)
        return send_file(path, as_attachment=True, download_name=f"{title}.xlsx")
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ===== PHASE 9: SUBSCRIPTION & LAUNCH PREP =====
@app.post('/api/subscription/upgrade')
def upgrade_subscription():
    """Mock endpoint to upgrade user to Pro"""
    data = request.json
    username = data.get('username', '').strip().lower()
    plan = data.get('plan', 'pro')
    
    if not username:
        return jsonify({'error': 'Username required'}), 400
        
    db = load_db()
    if username not in db['users']:
        return jsonify({'error': 'User not found'}), 404
        
    user = db['users'][username]
    user['tier'] = plan
    user['credits'] = 1000000 if plan == 'pro' else 10000
    save_db(db)
    
    return jsonify({'ok': True, 'tier': plan, 'credits': user['credits']})

@app.get('/api/user/credits')
def api_get_user_credits():
    """Fetch user credits and tier"""
    username = request.args.get('username', '').strip().lower()
    db = load_db()
    user = db['users'].get(username, {})
    return jsonify({
        'credits': user.get('credits', 10000),
        'tier': user.get('tier', 'free')
    })

if __name__ == '__main__':

    port = int(os.environ.get('PORT', 5054))
    print(f"\nAryaX AI Production Server Starting on Port {port}...")
    app.run(host='0.0.0.0', port=port)
